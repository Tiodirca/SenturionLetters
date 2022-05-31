from pptx import Presentation
from threading import Thread

estrofe_slide = ""
estrofe_slide_secundaria = ""


# metodo para gerar e salvar o slides com a letra

def adicionar_slides(prs, titulo, estrofe):
    blank_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_slide_layout)
    title = slide.shapes.title
    subtitulo = slide.placeholders[1]
    title.text = titulo
    subtitulo.text = estrofe


def gerar_slides(textos, titulo, tipo_modelo, caminho_salvar):
    # metodo para criar arquivo de slides
    class Thead_gerar_arquivo(Thread):  # instanciando uma nova thead para ser executada as instrucoes abaixo
        def __init__(self):
            Thread.__init__(self)

        def run(self):
            global estrofe_slide, estrofe_slide_secundaria
            if 'modelo_geral' in tipo_modelo:
                # passando o caminho do modelo que sera utilizado para gerar os slides
                caminho = "src/modelos_slides/modelo_geral.pptx"
                prs = Presentation(caminho)
            else:
                # passando o caminho do modelo que sera utilizado para gerar os slides
                caminho = "src/modelos_slides/modelo_geracao_fire.pptx"
                prs = Presentation(caminho)

            for estrofe in textos:
                # divindo a lista na quantidade de vezes passada como parametro para o split
                dividir_estrofe = str(estrofe).split("\n", 2)
                if len(dividir_estrofe) == 1:
                    estrofe_slide = dividir_estrofe[0]
                elif len(dividir_estrofe) >= 2:
                    # definindo que a variavel vai receber o valor pegando o primeiro index e o segundo index
                    estrofe_slide = dividir_estrofe[0] + "\n" + dividir_estrofe[1]
                adicionar_slides(prs, titulo, estrofe_slide)
                # verificando se a variavel contem mais de dois index
                if len(dividir_estrofe) > 2:
                    # divindo a lista a cada dois index e defindindo que a variavel vai receber o valor dividido
                    dividir_estrofe_secundaria = str(dividir_estrofe[2]).split("\n", 2)
                    # verificando se a variavel atende quais situacoes
                    if len(dividir_estrofe_secundaria) == 1:
                        estrofe_slide_secundaria = dividir_estrofe_secundaria[0]
                        adicionar_slides(prs, titulo, estrofe_slide_secundaria)

                    elif len(dividir_estrofe_secundaria) >= 2:
                        # definindo que a variavel vai receber um valor pegando o primeiro index e o segundo index
                        estrofe_slide_secundaria = dividir_estrofe_secundaria[0] + "\n" + dividir_estrofe_secundaria[1]
                        adicionar_slides(prs, titulo, estrofe_slide_secundaria)
                        # verificando se a variavel atende a seguinte situacao
                        if len(dividir_estrofe_secundaria) >= 3:
                            adicionar_slides(prs, titulo, dividir_estrofe_secundaria[2])

            # salvando o slide formato a ser salvo
            # ja passado na variavel caminho salvar
            prs.save(caminho_salvar)

    iniciar_thead = Thead_gerar_arquivo()
    iniciar_thead.start()
    return True
